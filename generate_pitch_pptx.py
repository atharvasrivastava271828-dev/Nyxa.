import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

# --- HUMANIZED ELEGANT DESIGN SYSTEM ---
COLOR_BG = RGBColor(0, 0, 0)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_GRAY_LIGHT = RGBColor(180, 180, 180)
COLOR_GRAY_DARK = RGBColor(100, 100, 100)
COLOR_GOLD = RGBColor(255, 215, 0)

FONT_HEADING = "Inter"
FONT_BODY = "Inter"

def set_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_slide_transition(slide, transition_type="fade"):
    sld = slide.element
    for child in list(sld):
        if child.tag.endswith('transition'):
            sld.remove(child)
            
    xml = f'<p:transition {nsdecls("p")} spd="slow"><p:push dir="u"/></p:transition>'
    transition_elem = parse_xml(xml)
    
    clrMapOvr = sld.find('{http://schemas.openxmlformats.org/presentationml/2006/main}clrMapOvr')
    timing = sld.find('{http://schemas.openxmlformats.org/presentationml/2006/main}timing')
    
    if timing is not None:
        timing.addprevious(transition_elem)
    elif clrMapOvr is not None:
        clrMapOvr.addnext(transition_elem)
    else:
        cSld = sld.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
        if cSld is not None:
            cSld.addnext(transition_elem)
        else:
            sld.append(transition_elem)

def add_line(slide, x, y, width, height=0):
    shape = slide.shapes.add_shape(
        9, # msoShapeLine
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    shape.line.color.rgb = COLOR_GRAY_DARK
    shape.line.width = Pt(1)
    return shape

def add_text(slide, text, x, y, w, h, font_size=20, font_color=COLOR_GRAY_LIGHT, bold=False, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = FONT_BODY if font_size < 30 else FONT_HEADING
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = align
    return tb

def create_humanized_deck(base_dir, yinyang_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # --- SLIDE 1: TITLE ---
    s1 = prs.slides.add_slide(blank)
    set_bg(s1)
    add_slide_transition(s1)
    
    if os.path.exists(yinyang_path):
        s1.shapes.add_picture(yinyang_path, Inches(6.066), Inches(0.8), height=Inches(1.2))
    
    add_text(s1, "Nyxa.", 1.0, 2.1, 11.333, 1.4, font_size=84, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_line(s1, 4.0, 3.6, 5.333)
    add_text(s1, "GET YOUR DAILY TASKS DONE IN SECONDS", 1.0, 3.9, 11.333, 0.5, font_size=20, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s1, "Simple, free, and private utilities designed for students, freelancers, and small businesses.", 1.0, 4.5, 11.333, 0.5, font_size=18, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)
    add_text(s1, "MAKING SURE NO ONE IS LEFT OUT", 1.0, 5.8, 11.333, 0.4, font_size=14, font_color=COLOR_GOLD, bold=True, align=PP_ALIGN.CENTER)

    # --- SLIDE 2: HYPOTHESIS & BUSINESS IDEAS ---
    s2 = prs.slides.add_slide(blank)
    set_bg(s2)
    add_slide_transition(s2)
    add_text(s2, "People don't want complex software.\nThey just want the result.", 1.0, 1.1, 11.333, 1.3, font_size=38, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s2, "When you need a GST invoice or a study quiz, you shouldn't have to learn a new app.", 1.0, 2.4, 11.333, 0.5, font_size=18, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)
    add_line(s2, 4.0, 3.0, 5.333)
    
    add_text(s2, "The Friction of Tools", 1.0, 3.5, 3.5, 0.4, font_size=20, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s2, "Current tools force you to create accounts, manage passwords, download apps, and pay monthly subscriptions just to do one simple task.", 1.0, 4.0, 3.5, 2.5, font_size=15, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)

    add_text(s2, "The Privacy Problem", 4.916, 3.5, 3.5, 0.4, font_size=20, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s2, "Basic utilities constantly track users, sell personal data, and force you to upload sensitive files to their servers just to get a result.", 4.916, 4.0, 3.5, 2.5, font_size=15, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)

    add_text(s2, "The Need for Outcomes", 8.833, 3.5, 3.5, 0.4, font_size=20, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s2, "Nobody wakes up wanting to use software. People simply want the invoice, the QR code, or the calculated marks, instantly.", 8.833, 4.0, 3.5, 2.5, font_size=15, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)

    # --- SLIDE 3: SOLUTION ---
    s3 = prs.slides.add_slide(blank)
    set_bg(s3)
    add_slide_transition(s3)
    add_text(s3, "The Marketplace for Finished Work", 1.0, 1.1, 11.333, 1.2, font_size=42, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s3, "Nyxa gives you the exact tools you need. 1. Pick  →  2. Type  →  3. Copy Result", 1.0, 2.3, 11.333, 0.5, font_size=18, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)
    add_line(s3, 4.0, 2.9, 5.333)

    add_text(s3, "Works Instantly", 1.0, 3.5, 3.5, 0.4, font_size=20, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s3, "No signups, no paywalls, no downloads. Pick a tool, type your details, and everything updates in real-time in your browser.", 1.0, 4.0, 3.5, 2.5, font_size=15, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)

    add_text(s3, "100% Private & Safe", 4.916, 3.5, 3.5, 0.4, font_size=20, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s3, "Your personal details never leave your browser. All computations happen client-side. Zero tracking, zero data selling.", 4.916, 4.0, 3.5, 2.5, font_size=15, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)

    add_text(s3, "Custom Task Board", 8.833, 3.5, 3.5, 0.4, font_size=20, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s3, "If a tool doesn't exist, post a request on our TaskBidder board. Developers build it for you, and we hold the payment safely in escrow.", 8.833, 4.0, 3.5, 2.5, font_size=15, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)

    # --- SLIDE 4: MARKET SIZE ---
    s4 = prs.slides.add_slide(blank)
    set_bg(s4)
    add_slide_transition(s4)
    add_text(s4, "Serving the Everyday Economy", 1.0, 1.1, 11.333, 1.0, font_size=38, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s4, "Millions of students, freelancers, and small businesses need simple daily utilities.", 1.0, 2.1, 11.333, 0.5, font_size=18, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)
    add_text(s4, "1.5 Billion+", 1.0, 2.7, 11.333, 1.8, font_size=110, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(s4, "Global Freelancers & Students", 1.0, 4.7, 3.5, 0.4, font_size=18, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s4, "A massive, underserved audience that relies on disjointed, ad-heavy websites for daily calculators, converters, and formatters.", 1.0, 5.2, 3.5, 1.8, font_size=14, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)

    add_text(s4, "The Digital Services Market", 4.916, 4.7, 3.5, 0.4, font_size=18, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s4, "Small businesses spend billions annually hiring out simple digital tasks that can be completely automated by focused tools.", 4.916, 5.2, 3.5, 1.8, font_size=14, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)

    add_text(s4, "Nyxa's Market Opportunity", 8.833, 4.7, 3.5, 0.4, font_size=18, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s4, "By becoming the default starting page for \"getting things done\", we capture massive daily recurring traffic and custom task requests.", 8.833, 5.2, 3.5, 1.8, font_size=14, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)

    # --- SLIDE 5: PRODUCT DEMO (TRANSFORMATIONS) ---
    s5 = prs.slides.add_slide(blank)
    set_bg(s5)
    add_slide_transition(s5)
    add_text(s5, "The Transformation", 1.0, 1.1, 11.333, 1.0, font_size=38, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s5, "1-Click Tools. Instant Results.", 1.0, 2.1, 11.333, 0.5, font_size=18, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)
    
    # Transformation 1
    add_text(s5, "Raw Details", 1.0, 3.2, 5.0, 0.8, font_size=32, font_color=COLOR_GRAY_DARK, bold=True, align=PP_ALIGN.RIGHT)
    add_text(s5, "→", 6.1, 3.2, 1.0, 0.8, font_size=32, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s5, "GST Invoice PDF", 7.2, 3.2, 5.0, 0.8, font_size=32, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.LEFT)

    # Transformation 2
    add_text(s5, "Bank Details", 1.0, 4.2, 5.0, 0.8, font_size=32, font_color=COLOR_GRAY_DARK, bold=True, align=PP_ALIGN.RIGHT)
    add_text(s5, "→", 6.1, 4.2, 1.0, 0.8, font_size=32, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s5, "UPI Payment QR", 7.2, 4.2, 5.0, 0.8, font_size=32, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.LEFT)

    # Transformation 3
    add_text(s5, "Syllabus Data", 1.0, 5.2, 5.0, 0.8, font_size=32, font_color=COLOR_GRAY_DARK, bold=True, align=PP_ALIGN.RIGHT)
    add_text(s5, "→", 6.1, 5.2, 1.0, 0.8, font_size=32, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s5, "Marks Calculator", 7.2, 5.2, 5.0, 0.8, font_size=32, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.LEFT)

    # --- SLIDE 6: BUSINESS MODEL ---
    s6 = prs.slides.add_slide(blank)
    set_bg(s6)
    add_slide_transition(s6)
    add_text(s6, "How We Monetize the Outcomes", 1.0, 1.1, 11.333, 1.0, font_size=38, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s6, "Core utilities remain completely free. We monetize custom requests and creator discovery.", 1.0, 2.1, 11.333, 0.5, font_size=18, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)

    add_text(s6, "15%", 1.0, 3.0, 5.3, 1.6, font_size=90, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s6, "TaskBidder Commission", 1.0, 4.7, 5.3, 0.4, font_size=22, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s6, "When a user requests a custom task or workflow to be built on our board, we take a 15% platform escrow fee upon successful delivery.", 1.0, 5.2, 5.3, 1.8, font_size=15, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)

    add_text(s6, "Ads", 7.0, 3.0, 5.3, 1.6, font_size=90, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s6, "Sponsored Discovery", 7.0, 4.7, 5.3, 0.4, font_size=22, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s6, "As the library grows, verified creators can pay to promote their specialized tools and API services to our high-intent traffic.", 7.0, 5.2, 5.3, 1.8, font_size=15, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)

    # --- SLIDE 7: TRACTION & SUPPLY ---
    s7 = prs.slides.add_slide(blank)
    set_bg(s7)
    add_slide_transition(s7)
    add_text(s7, "Built for Real Outcomes", 1.0, 1.1, 11.333, 1.0, font_size=38, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s7, "Traction from day one through high-utility free tools.", 1.0, 2.1, 11.333, 0.5, font_size=18, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)
    add_line(s7, 4.0, 2.7, 5.333)

    add_text(s7, "500+", 1.0, 3.3, 3.5, 0.8, font_size=60, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s7, "Tasks Completed", 1.0, 4.2, 3.5, 0.4, font_size=20, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)
    
    add_text(s7, "120+", 4.916, 3.3, 3.5, 0.8, font_size=60, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s7, "Verified Creators", 4.916, 4.2, 3.5, 0.4, font_size=20, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)

    add_text(s7, "Growing", 8.833, 3.3, 3.5, 0.8, font_size=60, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s7, "Free Utility Library", 8.833, 4.2, 3.5, 0.4, font_size=20, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)

    # --- SLIDE 8: COMPETITION ---
    s8 = prs.slides.add_slide(blank)
    set_bg(s8)
    add_slide_transition(s8)
    add_text(s8, "We Sell the Finished Job, Not the Tool", 1.0, 1.1, 11.333, 1.0, font_size=38, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s8, "How Nyxa disrupts the current bloated software market.", 1.0, 2.1, 11.333, 0.5, font_size=18, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)
    add_line(s8, 4.0, 2.7, 5.333)

    add_text(s8, "TRADITIONAL SAAS", 1.0, 3.3, 3.5, 0.4, font_size=16, font_color=COLOR_GRAY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s8, "Bloated & Expensive", 1.0, 3.8, 3.5, 0.4, font_size=18, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s8, "Force users into $10/mo subscriptions just to use a simple invoice generator or file converter once a month.", 1.0, 4.3, 3.5, 2.2, font_size=14, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)

    add_text(s8, "UTILITY WEBSITES", 4.916, 3.3, 3.5, 0.4, font_size=16, font_color=COLOR_GRAY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s8, "Unsafe & Ad-Heavy", 4.916, 3.8, 3.5, 0.4, font_size=18, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s8, "Riddled with popups, malware, and privacy risks, often uploading your private financial data to unknown servers.", 4.916, 4.3, 3.5, 2.2, font_size=14, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)

    add_text(s8, "NYXA PLATFORM", 8.833, 3.3, 3.5, 0.4, font_size=16, font_color=COLOR_GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s8, "Instant, Free & Private", 8.833, 3.8, 3.5, 0.4, font_size=18, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s8, "A unified, beautiful hub of tools that run instantly in your browser. No signups, completely private, and free to use.", 8.833, 4.3, 3.5, 2.2, font_size=14, font_color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    # --- SLIDE 9: ROADMAP ---
    s9 = prs.slides.add_slide(blank)
    set_bg(s9)
    add_slide_transition(s9)
    add_text(s9, "The Roadmap to Complex Outcomes", 1.0, 1.1, 11.333, 1.0, font_size=38, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s9, "Starting simple, scaling to full automation.", 1.0, 2.1, 11.333, 0.5, font_size=18, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)
    add_line(s9, 4.0, 2.7, 5.333)

    add_text(s9, "1. Dominate Everyday Utilities", 1.0, 3.3, 5.3, 0.4, font_size=20, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s9, "We are rapidly expanding our library of free tools—taxes, invoices, study aids—to capture massive organic traffic and become the default starting page for students and freelancers.", 1.0, 3.8, 5.3, 2.2, font_size=15, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)

    add_text(s9, "2. Chaining Custom Workflows", 7.0, 3.3, 5.3, 0.4, font_size=20, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s9, "As the TaskBidder board grows, creators will combine simple utilities into complex, paid business workflows—chaining a Market Research tool directly into a Pitch Deck generator.", 7.0, 3.8, 5.3, 2.2, font_size=15, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)

    # --- SLIDE 10: CLOSING ---
    s10 = prs.slides.add_slide(blank)
    set_bg(s10)
    add_slide_transition(s10)
    add_text(s10, "Building the Marketplace for Outcomes", 1.0, 1.1, 11.333, 1.0, font_size=36, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s10, "We are a student team building the internet we want to use.", 1.0, 2.0, 11.333, 0.5, font_size=18, font_color=COLOR_GRAY_LIGHT, align=PP_ALIGN.CENTER)
    
    add_text(s10, "People don't want tools.", 1.0, 3.2, 11.333, 0.6, font_size=32, font_color=COLOR_GRAY_LIGHT, bold=True, align=PP_ALIGN.CENTER)
    add_text(s10, "People want outcomes.", 1.0, 4.0, 11.333, 0.6, font_size=42, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)

    if os.path.exists(yinyang_path):
        s10.shapes.add_picture(yinyang_path, Inches(6.066), Inches(5.3), height=Inches(0.9))
        add_text(s10, "FOR THE LIGHT", 1.0, 6.3, 11.333, 0.4, font_size=14, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    else:
        add_text(s10, "FOR THE LIGHT", 1.0, 6.0, 11.333, 0.4, font_size=14, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Save outputs safely
    targets = [
        os.path.join(base_dir, "Nyxa_Pitch_Deck.pptx"),
        os.path.join(base_dir, "Nyxa_Pitch_Deck_Humanized.pptx")
    ]
    
    saved = False
    for path in targets:
        try:
            prs.save(path)
            print(f"Successfully generated humanized PowerPoint presentation at: {path}")
            saved = True
            break
        except Exception as e:
            print(f"Could not write to {path}: {e}")

    if not saved:
        alt_path = os.path.join(base_dir, "Nyxa_Pitch_Deck_Final.pptx")
        prs.save(alt_path)
        print(f"Saved to alternative path: {alt_path}")

if __name__ == "__main__":
    base_dir = r"d:\The Smiling Corporation\Nyxa"
    yinyang_path = os.path.join(base_dir, "public", "yinyang.png")
    create_humanized_deck(base_dir, yinyang_path)
